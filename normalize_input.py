import os
import mimetypes

# Required libs:
# pip install pdfplumber python-docx python-pptx openai-whisper nltk pysrt

import pdfplumber
from docx import Document
from pptx import Presentation
import whisper
import nltk
import pysrt

nltk.download('punkt')
from nltk.tokenize import sent_tokenize


def chunk_text(text, max_tokens=500):
    sentences = sent_tokenize(text)
    chunks = []
    current = ""
    for sent in sentences:
        if len(current.split()) + len(sent.split()) < max_tokens:
            current += " " + sent
        else:
            chunks.append(current.strip())
            current = sent
    if current:
        chunks.append(current.strip())
    return chunks


def extract_pdf(filepath):
    with pdfplumber.open(filepath) as pdf:
        text = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
    chunks = chunk_text(text)
    return [{"source": filepath, "type": "pdf", "chunk_index": i, "section": f"Page Chunk {i}", "text": chunk} for i, chunk in enumerate(chunks)]


def extract_docx(filepath):
    doc = Document(filepath)
    text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    chunks = chunk_text(text)
    return [{"source": filepath, "type": "docx", "chunk_index": i, "section": f"Paragraph Chunk {i}", "text": chunk} for i, chunk in enumerate(chunks)]


def extract_pptx(filepath):
    pres = Presentation(filepath)
    all_text = []
    for i, slide in enumerate(pres.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        if slide_text:
            text = "\n".join(slide_text)
            chunks = chunk_text(text)
            for j, chunk in enumerate(chunks):
                all_text.append({"source": filepath, "type": "pptx", "chunk_index": f"{i}_{j}", "section": f"Slide {i}", "text": chunk})
    return all_text


def extract_text_file(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        text = f.read()
    chunks = chunk_text(text)
    return [{"source": filepath, "type": "text", "chunk_index": i, "section": f"Text Chunk {i}", "text": chunk} for i, chunk in enumerate(chunks)]


def extract_subtitles(filepath):
    subs = pysrt.open(filepath)
    text_blocks = [sub.text.replace('\n', ' ') for sub in subs]
    full_text = ' '.join(text_blocks)
    chunks = chunk_text(full_text)
    return [{"source": filepath, "type": "subtitle", "chunk_index": i, "section": f"Subtitle Chunk {i}", "text": chunk} for i, chunk in enumerate(chunks)]


def transcribe_audio(filepath):
    model = whisper.load_model("base")
    result = model.transcribe(filepath)
    text = result["text"]
    chunks = chunk_text(text)
    return [{"source": filepath, "type": "audio", "chunk_index": i, "section": f"Transcript Chunk {i}", "text": chunk} for i, chunk in enumerate(chunks)]


def normalize_input(filepath):
    file_type, _ = mimetypes.guess_type(filepath)
    ext = os.path.splitext(filepath)[1].lower()

    if file_type == 'application/pdf':
        return extract_pdf(filepath)

    elif ext in ['.docx']:
        return extract_docx(filepath)

    elif ext in ['.pptx']:
        return extract_pptx(filepath)

    elif ext in ['.mp4', '.mp3', '.wav', '.m4a']:
        return transcribe_audio(filepath)

    elif ext in ['.txt', '.md']:
        return extract_text_file(filepath)

    elif ext in ['.srt', '.vtt']:
        return extract_subtitles(filepath)

    else:
        raise ValueError(f"Unsupported file type: {filepath}")
    
if __name__ == "__main__":
    test_file = "Memory.pdf"  # Change this to any supported file type you want to try
    chunks = normalize_input(test_file)
    for i, chunk in enumerate(chunks):
        print(f"\n--- Chunk {i} ---")
        print(chunk["text"])